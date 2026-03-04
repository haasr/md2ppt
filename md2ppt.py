"""
slides_builder.py
-----------------
Converts a Markdown string into a .pptx file using python-pptx.

Markdown conventions
--------------------
  # Title                   → Title slide (title only)
  # Title: Subtitle         → Title slide (split on first colon)
  ## Slide Title            → Content slide title
  Paragraphs under ##       → Plain text body items (no bullet)
  - Bullet items under ##   → Bulleted list items
  1. Numbered items under ## → Hanging-indent numbered plain text

Usage (as a module)
-------------------
  builder = SlidesBuilder(markdown_text, "output.pptx")
  builder.build()

Usage (CLI)
-----------
  python slides_builder.py input.md output.pptx
"""

import re
import sys
from dataclasses import dataclass, field
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE


# ── Layout indices ─────────────────────────────────────────────────────────────
_LAYOUT_TITLE_SLIDE   = 0   # "Title Slide"       – big title + subtitle
_LAYOUT_TITLE_CONTENT = 1   # "Title and Content" – title + body placeholder


# ── Data classes ───────────────────────────────────────────────────────────────

@dataclass
class SlideItem:
    """One logical item inside a content slide body."""
    PLAIN    = "plain"
    BULLET   = "bullet"
    NUMBERED = "numbered"

    kind:   str
    text:   str
    number: Optional[int] = None   # only used for NUMBERED items


@dataclass
class SlideData:
    """Parsed representation of a single slide."""
    TITLE_SLIDE   = "title"
    CONTENT_SLIDE = "content"

    kind:     str
    title:    str
    subtitle: Optional[str]   = None   # title slides only
    items:    list            = field(default_factory=list)


# ── Main class ─────────────────────────────────────────────────────────────────

class SlidesBuilder:
    """
    Parse a Markdown string and write a .pptx file.

    Parameters
    ----------
    markdown : str
        Source Markdown text (see module docstring for conventions).
    output_path : str
        Destination file path, e.g. "slides.pptx".
    slide_width_inches : float
        Slide width (default 13.33 for 16:9 widescreen).
    slide_height_inches : float
        Slide height (default 7.5 for 16:9 widescreen).
    theme_colors : dict | None
        Optional mapping of Office theme color slot names to hex strings
        (without the leading #).  Any subset of slots may be supplied;
        unspecified slots are left unchanged.

        Standard slot names:
            dk1, lt1, dk2, lt2,
            accent1, accent2, accent3, accent4, accent5, accent6,
            hlink, folHlink

        Example — ETSU palette::
 
            theme_colors = {
                "dk1":     "000000",   # Body text — black
                "lt1":     "FFFFFF",   # Slide background — white
                "dk2":     "00053E",   # Deepest navy
                "lt2":     "FFC72C",   # Bright gold — decorative/line use only
                "accent1": "003865",   # ETSU Blue          contrast ~8.5:1 ✓
                "accent2": "005A9E",   # Medium blue        contrast ~6.1:1 ✓
                "accent3": "7A6000",   # Dark amber/bronze  contrast ~7.2:1 ✓ (warm without being yellow)
                "accent4": "00053E",   # Deep navy repeat   contrast ~16:1  ✓
                "accent5": "1A5276",   # Steel blue         contrast ~7.8:1 ✓
                "accent6": "4A4A4A",   # Charcoal           contrast ~9.7:1 ✓
                "hlink":   "003865",   # Hyperlinks in ETSU blue
            }
    """

    def __init__(
        self,
        markdown: str,
        output_path: str,
        slide_width_inches: float  = 13.33,
        slide_height_inches: float = 7.5,
        theme_colors: Optional[dict] = None,
    ):
        if not output_path.lower().endswith(".pptx"):
            output_path = output_path + ".pptx"

        self.markdown            = markdown
        self.output_path         = output_path
        self.slide_width_inches  = slide_width_inches
        self.slide_height_inches = slide_height_inches
        self.theme_colors        = theme_colors
        self._prs                = None

    # ── Public entry point ─────────────────────────────────────────────────────

    def build(self) -> None:
        """Parse the Markdown and write the .pptx file."""
        slides = self._parse_markdown()
        self._prs = self._init_presentation()
        if self.theme_colors:
            self._apply_theme_colors(self.theme_colors)
        for slide_data in slides:
            if slide_data.kind == SlideData.TITLE_SLIDE:
                self._add_title_slide(slide_data)
            else:
                self._add_content_slide(slide_data)
        self._prs.save(self.output_path)
        print(f"Presentation saved → {self.output_path}")

    # ── Markdown parsing ───────────────────────────────────────────────────────

    def _parse_markdown(self) -> list:
        """
        Walk through the Markdown line by line and build a list of SlideData.
        """
        slides  = []
        current = None

        for raw_line in self.markdown.splitlines():
            line = raw_line.rstrip()

            # ── # Heading → title slide ───────────────────────────────────────
            if re.match(r'^#\s+', line) and not re.match(r'^##', line):
                text = re.sub(r'^#\s+', '', line).strip()
                if ':' in text:
                    title, subtitle = text.split(':', 1)
                    current = SlideData(
                        kind     = SlideData.TITLE_SLIDE,
                        title    = title.strip(),
                        subtitle = subtitle.strip(),
                    )
                else:
                    current = SlideData(kind=SlideData.TITLE_SLIDE, title=text)
                slides.append(current)

            # ── ## Heading → content slide ────────────────────────────────────
            elif re.match(r'^##\s+', line):
                title   = re.sub(r'^##\s+', '', line).strip()
                current = SlideData(kind=SlideData.CONTENT_SLIDE, title=title)
                slides.append(current)

            # ── - Bullet item ─────────────────────────────────────────────────
            elif re.match(r'^-\s+', line):
                if current and current.kind == SlideData.CONTENT_SLIDE:
                    text = re.sub(r'^-\s+', '', line).strip()
                    current.items.append(SlideItem(kind=SlideItem.BULLET, text=text))

            # ── 1. Numbered item ──────────────────────────────────────────────
            elif re.match(r'^\d+\.\s+', line):
                if current and current.kind == SlideData.CONTENT_SLIDE:
                    m = re.match(r'^(\d+)\.\s+(.*)', line)
                    if m:
                        number = int(m.group(1))
                        text   = m.group(2).strip()
                        current.items.append(
                            SlideItem(kind=SlideItem.NUMBERED, text=text, number=number)
                        )

            # ── Plain paragraph ───────────────────────────────────────────────
            elif line.strip():
                if current and current.kind == SlideData.CONTENT_SLIDE:
                    current.items.append(
                        SlideItem(kind=SlideItem.PLAIN, text=line.strip())
                    )

        return slides

    # ── Presentation initialisation ────────────────────────────────────────────

    def _init_presentation(self) -> Presentation:
        prs = Presentation()
        prs.slide_width  = Inches(self.slide_width_inches)
        prs.slide_height = Inches(self.slide_height_inches)
        return prs

    # ── Slide creation ─────────────────────────────────────────────────────────

    def _add_title_slide(self, data: SlideData) -> None:
        layout = self._prs.slide_layouts[_LAYOUT_TITLE_SLIDE]
        slide  = self._prs.slides.add_slide(layout)

        title_shape    = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        self._scale_shape(title_shape,    width=11.86)
        self._scale_shape(subtitle_shape, width=10.33)

        title_shape.text = data.title
        if data.subtitle:
            subtitle_shape.text = data.subtitle

    def _add_content_slide(self, data: SlideData) -> None:
        layout  = self._prs.slide_layouts[_LAYOUT_TITLE_CONTENT]
        slide   = self._prs.slides.add_slide(layout)

        title_shape   = slide.shapes.title
        content_shape = slide.placeholders[1]

        self._scale_shape(title_shape)
        self._scale_shape(content_shape)

        title_shape.text = data.title
        self._align_text(title_shape, PP_ALIGN.LEFT)

        for item in data.items:
            if item.kind == SlideItem.BULLET:
                self._add_paragraph(content_shape, item.text)   # keep native bullet
            elif item.kind == SlideItem.PLAIN:
                p = self._add_paragraph(content_shape, item.text)
                self._remove_bullet(p)
            elif item.kind == SlideItem.NUMBERED:
                # "N.\u2002 Text"  ← en-space gives a slightly wider gap than a
                # regular space, mimicking PowerPoint's own numbered list style.
                display = f"{item.number}.\u2002 {item.text}"
                p = self._add_paragraph(content_shape, display)
                self._remove_bullet(p)
                self._apply_hanging_indent(p)

    # ── Theme helpers ─────────────────────────────────────────────────────────

    def _apply_theme_colors(self, color_map: dict) -> None:
        """
        Patch the slide master's theme color scheme with custom colors.

        python-pptx exposes the theme as a generic blob-based Part with no
        high-level API, so we parse the XML directly, mutate the clrScheme
        element, then write the serialized bytes back to the part.

        Each key in color_map must be a valid Office theme slot name; each
        value is a 6-character hex string without a leading '#'.  Only the
        supplied slots are modified — everything else is left as-is.

        Valid slot names:
            dk1, lt1, dk2, lt2,
            accent1, accent2, accent3, accent4, accent5, accent6,
            hlink, folHlink
        """
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        # The theme lives as a relationship off the slide master part
        master_part = self._prs.slide_master.part
        theme_part = next(
            rel.target_part
            for rel in master_part.rels.values()
            if rel.reltype.split('/')[-1] == 'theme'
        )

        tree = etree.fromstring(theme_part.blob)
        clr_scheme = tree.find(f'.//{{{ns}}}clrScheme')
        if clr_scheme is None:
            return

        for slot_name, hex_color in color_map.items():
            slot = clr_scheme.find(f'{{{ns}}}{slot_name}')
            if slot is None:
                continue
            # Replace whatever child is there (sysClr or srgbClr) with srgbClr
            for child in list(slot):
                slot.remove(child)
            etree.SubElement(slot, f'{{{ns}}}srgbClr', val=hex_color)

        # Write the modified XML back to the part blob
        theme_part._blob = etree.tostring(
            tree, xml_declaration=True, encoding='UTF-8', standalone=True
        )

    # ── Shape / text-frame helpers ─────────────────────────────────────────────


    def _scale_shape(
        self,
        shape,
        width:        float = 12.33,
        margin_left:  int   = 0,
        margin_right: int   = 10,
        autosize:     bool  = False,
    ) -> None:
        """
        Resize a placeholder to the given width (inches) while preserving its
        original left / top / height so the slide layout isn't disturbed.
        """
        left   = shape.left
        top    = shape.top
        height = shape.height

        shape.width                   = Inches(width)
        shape.text_frame.margin_left  = margin_left
        shape.text_frame.margin_right = margin_right
        shape.left                    = left
        shape.top                     = top
        shape.height                  = height

        if not autosize:
            shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE

    @staticmethod
    def _align_text(shape, alignment=PP_ALIGN.LEFT) -> None:
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = alignment

    @staticmethod
    def _add_paragraph(shape, text: str, level: int = 0, bold: bool = False, size=None):
        """
        Append a paragraph to a shape's text frame, reusing the first empty
        paragraph that python-pptx always creates on a fresh placeholder.
        """
        tf = shape.text_frame
        tf.word_wrap = True

        if len(tf.paragraphs[0].text) == 0:
            paragraph = tf.paragraphs[0]
        else:
            paragraph = tf.add_paragraph()

        paragraph.text  = text
        paragraph.level = level

        if bold:
            paragraph.font.bold = True
        if size is not None:
            paragraph.font.size = Pt(size)

        return paragraph

    @staticmethod
    def _remove_bullet(paragraph) -> None:
        """
        Strip all bullet / list formatting from a paragraph by injecting
        <a:buNone/> into the paragraph-properties XML element.

        Also explicitly zeros out marL and indent so the paragraph does not
        inherit hanging-indent values from the slide layout's default lstStyle.
        Without this, plain-text paragraphs can appear indented due to layout
        defaults bleeding through even when no indent was set on the paragraph.
        """
        ns   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        pPr  = paragraph._p.get_or_add_pPr()
        tags = [
            'buChar', 'buAutoNum', 'buFont', 'buClr',
            'buSzPct', 'buSzPts', 'buSzTx', 'buNone',
        ]
        for tag in tags:
            existing = pPr.find(f'{{{ns}}}{tag}')
            if existing is not None:
                pPr.remove(existing)
        etree.SubElement(pPr, f'{{{ns}}}buNone')
        # Explicitly zero indentation so layout defaults cannot bleed through
        pPr.set('marL',   '0')
        pPr.set('indent', '0')

    @staticmethod
    def _apply_hanging_indent(paragraph, indent_inches: float = 0.64) -> None:
        """
        Give a paragraph a hanging indent so wrapped lines align under the
        text rather than the number label.

        marL  = total left indent from the text-box edge
        indent = negative value pulls the first line back to the left
        Net effect: first line starts at 0, continuation lines indent by
        `indent_inches`.
        """
        indent_emu = Inches(indent_inches)
        pPr = paragraph._p.get_or_add_pPr()
        pPr.set('marL',   str(indent_emu))
        pPr.set('indent', str(-indent_emu))


# ── CLI entry point ────────────────────────────────────────────────────────────

if __name__ == '__main__':
    if len(sys.argv) != 3:
        print("Usage: python slides_builder.py <input.md> <output.pptx>")
        sys.exit(1)

    md_path, out_path = sys.argv[1], sys.argv[2]
    with open(md_path, 'r', encoding='utf-8') as fh:
        markdown_text = fh.read()

    etsu_colors = {
        "dk1":     "000000",   # Body text — black
        "lt1":     "FFFFFF",   # Slide background — white
        "dk2":     "00053E",   # Deepest navy
        "lt2":     "FFC72C",   # Bright gold — decorative/line use only
        "accent1": "003865",   # ETSU Blue          contrast ~8.5:1 ✓
        "accent2": "005A9E",   # Medium blue        contrast ~6.1:1 ✓
        "accent3": "7A6000",   # Dark amber/bronze  contrast ~7.2:1 ✓ (warm without being yellow)
        "accent4": "00053E",   # Deep navy repeat   contrast ~16:1  ✓
        "accent5": "1A5276",   # Steel blue         contrast ~7.8:1 ✓
        "accent6": "4A4A4A",   # Charcoal           contrast ~9.7:1 ✓
        "hlink":   "003865",   # Hyperlinks in ETSU blue
    }

    SlidesBuilder(markdown_text, out_path, theme_colors=etsu_colors).build()